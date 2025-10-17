import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Output path
OUTPUT = os.path.join(os.path.dirname(__file__), 'wellbot_architecture.pptx')

prs = Presentation()
slide_layout = prs.slide_layouts[6]  # blank
slide = prs.slides.add_slide(slide_layout)

# Title
left = top = Inches(0.3)
width = Inches(9.0)
height = Inches(0.8)
shape = slide.shapes.add_textbox(left, top, width, height)
tf = shape.text_frame
tf.text = "WellBot – Main Architecture"
tf.paragraphs[0].font.size = Pt(32)

# Helper to add a labeled box

def add_box(x, y, w, h, text, fill=(0xFF, 0xF2, 0xCC), outline=(0x66, 0x66, 0x66), font_size=14, bold=False):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    fill_obj = box.fill
    fill_obj.solid()
    fill_obj.fore_color.rgb = RGBColor(*fill)
    line = box.line
    line.color.rgb = RGBColor(*outline)
    tf = box.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = RGBColor(0x00, 0x00, 0x00)
    return box

# Coordinates layout (rough grid)
# Row 1: User -> Frontend -> Backend
add_box(0.2, 1.3, 1.6, 0.9, "User\nBrowser", fill=(0xDA,0xE8,0xFC))
add_box(2.1, 1.3, 2.4, 1.0, "Streamlit Frontend\nfrontend/app1.py", fill=(0xFF,0xF2,0xCC), bold=True)
add_box(5.0, 1.3, 2.2, 1.0, "Backend API\nbackend/main.py", fill=(0xF8,0xCE,0xCC))

# Row 2: NLU / Actions / KB
add_box(5.0, 2.6, 2.2, 1.0, "Dialogue & NLU\nnlu.py / dialogue_manager.py", fill=(0xE1,0xD5,0xE7))
add_box(7.3, 2.6, 2.2, 1.0, "Actions Layer\nactions/*.py", fill=(0xD5,0xE8,0xD4))
add_box(9.6, 2.6, 2.4, 1.0, "Knowledge Base\nJSON / DB", fill=(0xFF,0xE6,0xCC))

# Row 3: Feedback / Analytics / Localization
add_box(2.1, 3.9, 2.4, 1.0, "Feedback API\n/feedback/*", fill=(0xCC,0xE5,0xFF))
add_box(4.6, 3.9, 2.4, 1.0, "Analytics\n/analytics/*", fill=(0xCD,0xEB,0x8B))
add_box(7.1, 3.9, 2.4, 1.0, "Localization\ntranslate() / localize_dates", fill=(0xF5,0xF5,0xF5))

# Simple connectors using lines (optional – minimal for clarity)
from pptx.enum.dml import MSO_LINE_DASH_STYLE

def add_connector(x1,y1,x2,y2):
    line = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.width = Pt(1.5)
    line.line.dash_style = None

# Main flow
add_connector(1.8,1.8,2.1,1.8)  # User -> Frontend
add_connector(4.5,1.8,5.0,1.8)  # Frontend -> Backend
add_connector(6.1,2.3,6.1,2.6)  # Backend down to NLU
add_connector(7.2,3.1,7.3,3.1)  # NLU to Actions
add_connector(9.5,3.1,9.6,3.1)  # Actions to KB
add_connector(3.3,2.3,3.3,3.9)  # Frontend to Feedback API
add_connector(5.7,2.3,5.7,3.9)  # Backend to Analytics
add_connector(6.8,2.3,8.3,3.9)  # Backend to Localization (diagonal approximation)

# Save
prs.save(OUTPUT)
print(f"Generated: {OUTPUT}")
